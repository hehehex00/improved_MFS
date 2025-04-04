"""PPTX File Handler."""
from pptx import Presentation

from data_toolbox.multi_file_search.utils import (
    build_result,
    match_function,
)


def search_pptx(file, search_terms, search_options):
    """Search PPTX for Search Terms.

    Args:
    ----
        file (file): a PPTX file uploaded through streamlit's UI
        search_terms (list): Keywords to search the file for
        search_options (dictionary): configuration for search

    Returns:
    -------
        list: search results

    """
    results = []

    # Read the file
    try:
        ppt_reader = Presentation(file)
    except Exception:  # noqa: BLE001
        return [{
            "file": file.name,
            "location": "Error reading file",
        }]

    # Search each slide:
    for slide_number, slide in enumerate(ppt_reader.slides):
        slide_content = []
        # extract text from slide:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_content.append(shape.text)  # noqa: PERF401 (ignored for readability)

        for line in slide_content:
            matched_terms_in_line = []
            # check each search term:
            for term in search_terms:
                if match_function(line, term, search_options):
                    matched_terms_in_line.append(term)  # noqa: PERF401 (ignored for readability)
            # if one or more search term was found...
            if len(matched_terms_in_line):
                # ...add to the results
                results.append(
                    build_result(
                        file_name=file.name,
                        location=f"Slide {slide_number + 1}",
                        search_terms=matched_terms_in_line,
                        original_content=line,
                        location_context="",
                    ),
                )
    return results
